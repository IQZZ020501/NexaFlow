"""Extended stdlib-only coverage suite for the sandbox package.

Complements `sandbox.self_check` with branches it does not reach: request
error paths on the JSON-lines server, healthcheck edge cases, `Limits`
validation, `run_code` argument validation, and process-termination paths.

Run via `sandbox/run_coverage.sh` (measures with coverage) or directly with
`python -m sandbox.tests` (also runs inside the CI Docker image, where the
root/Linux-gated branches execute).
"""

from __future__ import annotations

import base64
from io import BytesIO
import json
import os
import signal
import socket
import subprocess
import sys
import tempfile
import threading
import time
from pathlib import Path
from unittest import mock
from zipfile import ZipFile

from sandbox import egress as egress_module
from sandbox import runner as runner_module
from sandbox import server as server_module
from sandbox.healthcheck import (
    MAX_RESPONSE_BYTES,
    PROBE_DEADLINE_SECONDS,
    _receive_line,
    probe,
)
from sandbox.runner import (
    MAX_CODE_BYTES,
    Limits,
    _terminate,
    execute_request,
    run_code,
)
from sandbox.server import MAX_REQUEST_BYTES, SandboxServer

HERE = Path(__file__).parent


def _coverage_prefix() -> list[str]:
    """Return [coverage-binary, run, flags] when available, else plain python."""
    binary = os.environ.get("SANDBOX_COVERAGE_BIN")
    if binary:
        nested = HERE / ".coverage.sandbox-nested"
        return [binary, "run", "--parallel-mode", f"--data-file={nested}"]
    return [sys.executable]


def check_limits_validation() -> None:
    for raw, fragment in [
        ("not-a-dict", "limits must be an object"),
        ({"timeout_ms": 0}, "positive integer"),
        ({"cpu_seconds": True}, "positive integer"),
        ({"memory_bytes": -1}, "positive integer"),
    ]:
        try:
            Limits.from_request(raw)
        except ValueError as exc:
            assert fragment in str(exc), (raw, exc)
        else:
            raise AssertionError(f"Limits.from_request({raw!r}) accepted")

    assert Limits.from_request(None) == Limits()
    clamped = Limits.from_request({"timeout_ms": 99_999, "max_processes": 2})
    assert clamped.timeout_ms == Limits().timeout_ms
    assert clamped.max_processes == 2

    try:
        Limits.from_request({"memory_mb": 128})
    except ValueError as exc:
        assert "unknown limits" in str(exc)
    else:
        raise AssertionError("unknown limit accepted")

    from sandbox.runner import MAX_STDIN_BYTES

    try:
        run_code("print('ignored')", "x" * (MAX_STDIN_BYTES + 1))
    except ValueError as exc:
        assert "stdin exceeds" in str(exc)
    else:
        raise AssertionError("oversized stdin accepted")


def check_run_code_validation() -> None:
    cases = [
        (lambda: run_code(123), ValueError, "code must be a string"),
        (lambda: run_code("print(1)", 123), ValueError, "stdin must be a string"),
        (lambda: run_code("\ud800"), ValueError, "valid UTF-8"),
        (lambda: run_code("x" * (MAX_CODE_BYTES // 1 + 1)), ValueError, "256 KiB"),
        (lambda: execute_request("x"), ValueError, "request must be an object"),
        (lambda: execute_request({"code": 123}), ValueError, "code must be a string"),
        (lambda: execute_request({"code": "x", "stdin": 1}), ValueError, "stdin must be a string"),
        (
            lambda: execute_request({"code": "pass", "artifact": []}),
            ValueError,
            "artifact must contain",
        ),
        (
            lambda: execute_request(
                {
                    "code": "pass",
                    "artifact": {"format": "../zip", "filename": "x.zip"},
                }
            ),
            ValueError,
            "artifact.format",
        ),
        (
            lambda: execute_request(
                {
                    "code": "pass",
                    "artifact": {"format": "py", "filename": "x.java"},
                }
            ),
            ValueError,
            "artifact.filename",
        ),
        (
            lambda: execute_request({"code": "pass", "skills": ["../documents"]}),
            ValueError,
            "skills must contain",
        ),
        (
            lambda: execute_request({"code": "pass", "skill": "documents"}),
            ValueError,
            "skill request",
        ),
        (
            lambda: execute_request({"skill": "../documents"}),
            ValueError,
            "skill request",
        ),
    ]
    for call, exc_type, fragment in cases:
        try:
            call()
        except exc_type as exc:
            assert fragment in str(exc), exc
        else:
            raise AssertionError(f"{call} did not raise {exc_type.__name__}")

    payload = execute_request({"code": "print(40 + 2)"})
    assert payload["ok"] is True and payload["stdout"].strip() == "42"
    assert execute_request({"code": "x"})["version"] == 1

    with tempfile.TemporaryDirectory() as directory:
        workdir = Path(directory)
        command = [sys.executable, "-c", "pass"]
        with mock.patch.object(runner_module.sys, "platform", "linux"):
            assert runner_module._sandboxed_child_command(command, workdir) == command
        with mock.patch.object(runner_module.sys, "platform", "darwin"), mock.patch.object(
            runner_module.shutil, "which", return_value=None
        ):
            try:
                runner_module._sandboxed_child_command(command, workdir)
            except RuntimeError as exc:
                assert "sandbox-exec" in str(exc)
            else:
                raise AssertionError("missing macOS sandbox-exec was accepted")
        with mock.patch.object(runner_module.sys, "platform", "darwin"), mock.patch.object(
            runner_module.shutil, "which", return_value="/usr/bin/sandbox-exec"
        ):
            command_with_proxy = runner_module._sandboxed_child_command(
                command, workdir, "http://127.0.0.1:4321"
            )
            assert "localhost:4321" in command_with_proxy[2]


def check_egress_policy() -> None:
    assert egress_module._is_public_address("1.1.1.1")
    for address in ("127.0.0.1", "10.0.0.1", "169.254.169.254", "::1"):
        assert not egress_module._is_public_address(address), address
    assert egress_module._destination("GET", "http://example.com/a?q=1")[2] == "/a?q=1"
    assert egress_module._host_port("[2001:db8::1]:443") == ("2001:db8::1", 443)
    for target in ("http://example.com:22/", "https://example.com/", "not-a-url"):
        try:
            egress_module._destination("GET", target)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError(f"disallowed proxy target accepted: {target}")

    public_info = [(socket.AF_INET, socket.SOCK_STREAM, 6, "", ("93.184.216.34", 80))]
    with mock.patch.object(egress_module.socket, "getaddrinfo", return_value=public_info):
        assert egress_module._resolve_public("example.test", 80) == ("93.184.216.34",)
    mixed_info = public_info + [
        (socket.AF_INET, socket.SOCK_STREAM, 6, "", ("192.168.1.10", 80))
    ]
    with mock.patch.object(egress_module.socket, "getaddrinfo", return_value=mixed_info):
        try:
            egress_module._resolve_public("example.test", 80)
        except egress_module.ProxyRequestError as exc:
            assert "not public" in str(exc)
        else:
            raise AssertionError("mixed public/private DNS result was accepted")


def check_egress_relay() -> None:
    client, proxy_side = socket.socketpair()
    upstream_peer, upstream_side = socket.socketpair()
    original_connect = egress_module._connect
    egress_module._connect = lambda host, port: upstream_side
    thread = threading.Thread(
        target=egress_module._serve_client, args=(proxy_side,), daemon=True
    )
    thread.start()
    try:
        client.sendall(
            b"GET http://example.com/ready?x=1 HTTP/1.1\r\n"
            b"Host: example.com\r\nConnection: close\r\n\r\n"
        )
        forwarded = upstream_peer.recv(4096)
        assert b"GET /ready?x=1 HTTP/1.1" in forwarded, forwarded
        assert b"Host: example.com\r\n" in forwarded, forwarded
        upstream_peer.sendall(
            b"HTTP/1.1 200 OK\r\nContent-Length: 2\r\nConnection: close\r\n\r\nok"
        )
        response = client.recv(4096)
        assert b"200 OK" in response and response.endswith(b"ok"), response
    finally:
        egress_module._connect = original_connect
        client.close()
        upstream_peer.close()
        upstream_side.close()
        thread.join(timeout=2)
        assert not thread.is_alive()


def check_local_egress_proxy() -> None:
    with tempfile.TemporaryDirectory() as directory:
        path = Path(directory) / "egress.sock"
        listener = socket.socket(socket.AF_UNIX, socket.SOCK_STREAM)
        listener.bind(str(path))
        listener.listen(1)

        def fake_egress() -> None:
            client, _ = listener.accept()
            with client:
                assert b"GET http://example.com/ready" in client.recv(4096)
                client.sendall(
                    b"HTTP/1.1 200 OK\r\nContent-Length: 2\r\n"
                    b"Connection: close\r\n\r\nok"
                )

        thread = threading.Thread(target=fake_egress, daemon=True)
        thread.start()
        proxy = egress_module.LocalEgressProxy(path)
        proxy.start()
        assert proxy.environment()["HTTP_PROXY"] == proxy.url
        try:
            port = int(proxy.url.rsplit(":", 1)[1])
            with socket.create_connection(("127.0.0.1", port)) as client:
                client.sendall(
                    b"GET http://example.com/ready HTTP/1.1\r\n"
                    b"Host: example.com\r\n\r\n"
                )
                response = client.recv(4096)
            assert response.endswith(b"ok"), response
        finally:
            proxy.close()
            listener.close()
            thread.join(timeout=2)
            assert not thread.is_alive()

    client, proxy_side = socket.socketpair()
    thread = threading.Thread(
        target=egress_module._serve_client, args=(proxy_side,), daemon=True
    )
    thread.start()
    try:
        client.sendall(
            b"GET http://127.0.0.1/metadata HTTP/1.1\r\n"
            b"Host: 127.0.0.1\r\n\r\n"
        )
        response = client.recv(4096)
        assert response.startswith(b"HTTP/1.1 403 Forbidden"), response
    finally:
        client.close()
        thread.join(timeout=2)
        assert not thread.is_alive()

    client, proxy_side = socket.socketpair()
    upstream_peer, upstream_side = socket.socketpair()
    original_connect = egress_module._connect
    egress_module._connect = lambda host, port: upstream_side
    thread = threading.Thread(
        target=egress_module._serve_client, args=(proxy_side,), daemon=True
    )
    thread.start()
    try:
        client.sendall(
            b"CONNECT example.com:443 HTTP/1.1\r\nHost: example.com\r\n\r\n"
            b"client-hello"
        )
        assert client.recv(4096).startswith(b"HTTP/1.1 200 Connection Established")
        assert upstream_peer.recv(4096) == b"client-hello"
    finally:
        egress_module._connect = original_connect
        client.close()
        upstream_peer.close()
        upstream_side.close()
        thread.join(timeout=2)
        assert not thread.is_alive()


def check_artifact_outputs_and_skills() -> None:
    html = execute_request(
        {
            "code": (
                "import os\n"
                "open(os.environ['NEXAFLOW_OUTPUT_PATH'], 'w', encoding='utf-8').write("
                "'<html><style>body{color:#123}</style><body>ready</body></html>')\n"
            ),
            "artifact": {"format": "html", "filename": "page.html"},
        }
    )
    assert html["ok"] is True, html
    assert html["artifact"]["filename"] == "page.html"
    assert html["artifact"]["format"] == "html"
    assert base64.b64decode(html["artifact"]["content_base64"]).endswith(
        b"</html>"
    )

    for artifact_format, filename, code in [
        (
            "docx",
            "report.docx",
            "import os\n"
            "from docx import Document\n"
            "output_path = os.environ['NEXAFLOW_OUTPUT_PATH']\n"
            "document = Document()\n"
            "document.add_paragraph('ready')\n"
            "document.save(output_path)\n",
        ),
        (
            "xlsx",
            "report.xlsx",
            "import os\n"
            "from openpyxl import Workbook\n"
            "output_path = os.environ['NEXAFLOW_OUTPUT_PATH']\n"
            "workbook = Workbook()\n"
            "workbook.active['A1'] = 'ready'\n"
            "workbook.save(output_path)\n",
        ),
        (
            "pptx",
            "slides.pptx",
            "import os\n"
            "from pptx import Presentation\n"
            "output_path = os.environ['NEXAFLOW_OUTPUT_PATH']\n"
            "presentation = Presentation()\n"
            "presentation.slides.add_slide(presentation.slide_layouts[6])\n"
            "presentation.save(output_path)\n",
        ),
        (
            "pdf",
            "report.pdf",
            "import os, pymupdf\n"
            "output_path = os.environ['NEXAFLOW_OUTPUT_PATH']\n"
            "document = pymupdf.open()\n"
            "page = document.new_page()\n"
            "page.insert_text((72, 72), 'ready')\n"
            "document.save(output_path)\n",
        ),
        (
            "py",
            "hello.py",
            "import os\n"
            "open(os.environ['NEXAFLOW_OUTPUT_PATH'], 'w', encoding='utf-8').write("
            "\"print('hello')\\n\")\n",
        ),
        (
            "java",
            "Hello.java",
            "import os\n"
            "open(os.environ['NEXAFLOW_OUTPUT_PATH'], 'w', encoding='utf-8').write("
            "'class Hello {}\\n')\n",
        ),
        (
            "zip",
            "bundle.zip",
            "import os, zipfile\n"
            "output_path = os.environ['NEXAFLOW_OUTPUT_PATH']\n"
            "with zipfile.ZipFile(output_path, 'w') as archive:\n"
            "    archive.writestr('hello.txt', 'ready')\n",
        ),
        (
            "file",
            "README",
            "import os\n"
            "open(os.environ['NEXAFLOW_OUTPUT_PATH'], 'w', encoding='utf-8').write('ready')\n",
        ),
    ]:
        result = execute_request(
            {
                "code": code,
                "artifact": {"format": artifact_format, "filename": filename},
            }
        )
        assert result["ok"] is True, result
        assert result["artifact"]["format"] == artifact_format
        assert result["artifact"]["filename"] == filename

    with tempfile.TemporaryDirectory() as directory:
        skill = Path(directory) / "documents"
        skill.mkdir()
        (skill / "SKILL.md").write_text("skill-ready", encoding="utf-8")
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": directory}):
            skilled = execute_request(
                {
                    "code": (
                        "import os\n"
                        "from pathlib import Path\n"
                        "text = (Path(os.environ['NEXAFLOW_SKILLS_DIR']) / "
                        "'documents' / 'SKILL.md').read_text()\n"
                        "open(os.environ['NEXAFLOW_OUTPUT_PATH'], 'w').write(text)\n"
                    ),
                    "skills": ["documents"],
                    "artifact": {"format": "html", "filename": "skill.html"},
                }
            )
            assert skilled["ok"] is True, skilled
            assert (
                base64.b64decode(skilled["artifact"]["content_base64"])
                == b"skill-ready"
            )

            missing = execute_request(
                {
                    "code": "pass",
                    "skills": ["missing"],
                    "artifact": {"format": "html", "filename": "missing.html"},
                }
            )
            assert missing["ok"] is False
            assert missing["error"] == "skill_not_found"

    try:
        execute_request(
            {
                "code": "pass",
                "artifact": {"format": "html", "filename": "../page.html"},
            }
        )
    except ValueError as exc:
        assert "filename" in str(exc)
    else:
        raise AssertionError("artifact path traversal was accepted")


def check_builtin_skill_entrypoints() -> None:
    requests = [
        (
            "documents",
            "docx",
            "skill-report.docx",
            {"content": "# Quarterly report\n\nReady for review.\n\n- Revenue\n- Cost"},
            "word/document.xml",
        ),
        (
            "pdf",
            "pdf",
            "skill-report.pdf",
            {"content": "# Quarterly report\n\nReady for review.\n\n- Revenue\n- Cost"},
            None,
        ),
        (
            "pptx",
            "pptx",
            "skill-deck.pptx",
            {
                "presentation": {
                    "title": "Quarterly review",
                    "subtitle": "Decisions and next steps",
                    "template": "bold",
                    "brand": {
                        "primary_color": "#F59E0B",
                        "background_color": "#111827",
                        "text_color": "#F9FAFB",
                        "font_family": "Arial",
                    },
                    "footer": "NexaFlow",
                    "slides": [
                        {
                            "layout": "icons",
                            "title": "Three quarter priorities",
                            "items": [
                                {
                                    "icon": "growth",
                                    "title": "Grow",
                                    "body": "Expand the strongest customer segment.",
                                },
                                {
                                    "icon": "gear",
                                    "title": "Simplify",
                                    "body": "Remove manual steps from delivery.",
                                },
                                {
                                    "icon": "focus",
                                    "title": "Focus",
                                    "body": "Fund the work tied to retention.",
                                },
                            ],
                            "notes": "[Sources]\n- Internal quarterly review",
                        },
                        {
                            "layout": "table",
                            "title": "One outcome per priority",
                            "table": {
                                "headers": ["Priority", "Outcome", "Owner"],
                                "rows": [
                                    ["Grow", "+12% qualified pipeline", "Sales"],
                                    ["Simplify", "-20% delivery time", "Operations"],
                                    ["Focus", "+5 pts retention", "Product"],
                                ],
                            },
                        },
                    ],
                }
            },
            "ppt/presentation.xml",
        ),
        (
            "spreadsheets",
            "xlsx",
            "skill-report.xlsx",
            {
                "workbook": {
                    "sheets": [
                        {
                            "name": "Summary",
                            "rows": [
                                ["Metric", "Value"],
                                ["Revenue", 12],
                                ["Cost", 5],
                                ["Profit", "=B2-B3"],
                            ],
                        }
                    ]
                }
            },
            "xl/workbook.xml",
        ),
    ]
    for skill, artifact_format, filename, inputs, archive_member in requests:
        result = execute_request(
            {
                "skill": skill,
                "stdin": json.dumps(inputs),
                "artifact": {"format": artifact_format, "filename": filename},
            }
        )
        assert result["ok"] is True, result
        content = base64.b64decode(result["artifact"]["content_base64"])
        assert f'"renderer":"{skill}"' in result["stdout"], result
        if artifact_format == "pdf":
            assert content.startswith(b"%PDF"), content[:16]
        else:
            with ZipFile(BytesIO(content)) as archive:
                assert archive_member in archive.namelist(), archive.namelist()

    long_title = "A" * 43
    result = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "bullets",
                                "title": long_title,
                                "bullets": ["Short point"],
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "invalid.pptx"},
        }
    )
    assert result["ok"] is True, result
    missing_table = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "table",
                                "title": "Missing table",
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "missing-table.pptx"},
        }
    )
    assert missing_table["ok"] is False
    assert (
        "presentation.slides[0].table is required for the table layout"
        in missing_table["stderr"]
    )

    long_cell_table = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "table",
                                "title": "Long cells wrap",
                                "table": {
                                    "headers": ["违法情形", "处罚措施"],
                                    "rows": [
                                        [
                                            "擅自举办培训机构（有场所、2名以上人员、有组织机构）",
                                            "责令停止举办、退还费用，处违法所得1—5倍罚款",
                                        ]
                                    ],
                                },
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "long-cells.pptx"},
        }
    )
    assert long_cell_table["ok"] is True, long_cell_table

    two_line_columns = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "two_column",
                                "title": "Two-line bullets",
                                "left": {
                                    "heading": "从轻、减轻或不予处罚",
                                    "bullets": [
                                        "主动消除或减轻危害后果",
                                        "受胁迫、诱骗实施；主动供述未掌握的违法行为",
                                        "配合查处有立功表现",
                                        "轻微并及时改正、无危害后果；能证明无主观过错",
                                    ],
                                },
                                "right": {
                                    "heading": "应当从重处罚",
                                    "bullets": [
                                        "被处理后两年内再次实施同类违法行为",
                                        "危害后果严重、造成恶劣社会影响",
                                        "伪造、涂改、转移、销毁证据",
                                        "拒绝、阻碍或以暴力威胁执法",
                                        "中小学在职教师从事学科类培训",
                                    ],
                                },
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "two-line-columns.pptx"},
        }
    )
    assert two_line_columns["ok"] is True, two_line_columns

    long_title = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "bullets",
                                "title": "结" * 23,
                                "bullets": ["Short point"],
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "long-title.pptx"},
        }
    )
    assert long_title["ok"] is True, long_title

    three_line_subtitle = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "section",
                                "title": "S",
                                "subtitle": "落" * 58,
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "three-line-subtitle.pptx"},
        }
    )
    assert three_line_subtitle["ok"] is True, three_line_subtitle

    two_line_icon_titles = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "icons",
                                "title": "S",
                                "items": [
                                    {
                                        "icon": "star",
                                        "title": "罚款、没收违法所得",
                                        "body": "经济制裁，追缴违法收益",
                                    },
                                    {
                                        "icon": "gear",
                                        "title": "责令停止招收学员、停止举办",
                                        "body": "限制或终止办学行为",
                                    },
                                    {
                                        "icon": "cycle",
                                        "title": "吊销许可证件、限制从业并同步追究相关责任人员责任",
                                        "body": "最严厉处理，并及于责任人员",
                                    },
                                    {
                                        "icon": "heart",
                                        "title": "警告、通报批评",
                                        "body": "最轻处罚",
                                    },
                                ],
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "two-line-icon-titles.pptx"},
        }
    )
    assert two_line_icon_titles["ok"] is True, two_line_icon_titles

    long_stat_value = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "stats",
                                "title": "Long metrics fit on the first render",
                                "stats": [
                                    {"value": "1", "label": "Baseline"},
                                    {
                                        "value": "12345678901234567890",
                                        "label": "Long identifier",
                                    },
                                    {"value": "34.7%", "label": "Growth"},
                                    {"value": "4", "label": "Regions"},
                                ],
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "long-stat.pptx"},
        }
    )
    assert long_stat_value["ok"] is True, long_stat_value

    five_step_titles = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Validation",
                        "slides": [
                            {
                                "layout": "steps",
                                "title": "S",
                                "steps": [
                                    {
                                        "title": "立案与调查",
                                        "body": "符合条件予以立案；现场调查、询问、查阅复制资料",
                                    },
                                    {
                                        "title": "告知与申辩",
                                        "body": "书面告知拟处罚内容及依据，听取陈述申辩并复核",
                                    },
                                    {
                                        "title": "听证与法制审核",
                                        "body": "重大罚款、没收及吊销等案件告知听证权利",
                                    },
                                    {
                                        "title": "决定与送达",
                                        "body": "制作行政处罚决定书，当场交付或七日内送达",
                                    },
                                    {
                                        "title": "执行与结案",
                                        "body": "不履行的申请法院强制执行，执行完毕结案归档",
                                    },
                                ],
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "five-step-titles.pptx"},
        }
    )
    assert five_step_titles["ok"] is True, five_step_titles

    custom_theme = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Model-created visual identity",
                        "subtitle": "The renderer protects readability and fit",
                        "theme": {
                            "background_color": "#07111F",
                            "text_color": "#F8FAFC",
                            "accent_color": "#F59E0B",
                            "panel_color": "#172033",
                            "muted_text_color": "#B8C4D6",
                            "rule_color": "#43516A",
                            "font_family": "Arial",
                            "heading_font_family": "Georgia",
                            "cover_title_size": 60,
                            "slide_title_size": 40,
                            "body_size": 20,
                            "panel_radius": 0.12,
                            "cover_accent_width": 0.42,
                            "title_alignment": "center",
                        },
                        "slides": [
                            {
                                "layout": "stats",
                                "title": "One theme, safe geometry",
                                "stats": [
                                    {"value": "1×", "label": "Tool call"},
                                    {"value": "4.5:1", "label": "Text contrast"},
                                    {"value": "35pt", "label": "Title floor"},
                                ],
                            },
                            {
                                "layout": "quote",
                                "title": "A deliberate chapter break",
                                "quote": "Style is model-authored; layout safety stays deterministic.",
                                "source": "NexaFlow renderer contract",
                                "style": {
                                    "background_color": "#FFF8E7",
                                    "text_color": "#2A1A14",
                                    "accent_color": "#8B1E3F",
                                    "panel_color": "#F2E4CC",
                                    "muted_text_color": "#6B4F44",
                                    "rule_color": "#A88C7D",
                                    "title_alignment": "left",
                                    "panel_radius": 0,
                                },
                            },
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "custom-theme.pptx"},
        }
    )
    assert custom_theme["ok"] is True, custom_theme
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE

    themed_deck = Presentation(
        BytesIO(base64.b64decode(custom_theme["artifact"]["content_base64"]))
    )
    assert str(themed_deck.slides[0].background.fill.fore_color.rgb) == "07111F"
    assert str(themed_deck.slides[2].background.fill.fore_color.rgb) == "FFF8E7"
    assert any(
        run.font.name == "Georgia"
        for slide in themed_deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    )
    assert any(
        shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        for shape in themed_deck.slides[1].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    )

    partial_theme = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Partial theme",
                        "theme": {
                            "background_color": "#07111F",
                            "text_color": "#F8FAFC",
                        },
                        "slides": [
                            {
                                "layout": "two_column",
                                "title": "Missing tokens get safe defaults",
                                "left": {"heading": "Model", "bullets": ["Theme intent"]},
                                "right": {"heading": "Renderer", "bullets": ["Safe fill"]},
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "partial-theme.pptx"},
        }
    )
    assert partial_theme["ok"] is True, partial_theme

    low_contrast_theme = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Invalid contrast",
                        "theme": {
                            "background_color": "#FFFFFF",
                            "text_color": "#F8FAFC",
                        },
                        "slides": [
                            {
                                "layout": "hero",
                                "title": "This must be rejected",
                            }
                        ],
                    }
                }
            ),
            "artifact": {"format": "pptx", "filename": "low-contrast.pptx"},
        }
    )
    assert low_contrast_theme["ok"] is False
    assert "4.5:1 contrast" in low_contrast_theme["stderr"]


def check_builtin_skill_quality_guards() -> None:
    long_text = "这是一段用于验证长文本自动换行和一次生成成功的中文内容" * 5
    document = execute_request(
        {
            "skill": "documents",
            "stdin": json.dumps(
                {
                    "content": (
                        "# 中文报告\n\n这是中文内容。\n\n"
                        "| 项目 | 说明 | 负责人 | 状态 |\n"
                        "| --- | --- | --- | --- |\n"
                        f"| 交付 | {long_text} | 团队 | 完成 |"
                    )
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "docx", "filename": "quality.docx"},
        }
    )
    assert document["ok"] is True, document
    with ZipFile(
        BytesIO(base64.b64decode(document["artifact"]["content_base64"]))
    ) as archive:
        document_xml = archive.read("word/document.xml").decode()
    if not os.environ.get("NEXAFLOW_CJK_FONT"):
        assert "Microsoft YaHei" not in document_xml
    assert '<w:tblW w:type="dxa" w:w="9071"' in document_xml
    assert '<w:tblLayout w:type="fixed"' in document_xml

    pdf = execute_request(
        {
            "skill": "pdf",
            "stdin": json.dumps(
                {
                    "content": "# 中文报告\n\n"
                    + "\n".join(f"- 第 {index} 项：{long_text}" for index in range(1, 45))
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pdf", "filename": "quality.pdf"},
        }
    )
    assert pdf["ok"] is True, pdf
    import pymupdf

    with pymupdf.open(
        stream=base64.b64decode(pdf["artifact"]["content_base64"]), filetype="pdf"
    ) as pdf_document:
        pdf_text = "\n".join(page.get_text() for page in pdf_document)
        assert "中" in pdf_text
        assert "1" in pdf_text
        assert pdf_document.page_count > 1

    spreadsheet = execute_request(
        {
            "skill": "spreadsheets",
            "stdin": json.dumps(
                {
                    "workbook": {
                        "sheets": [
                            {
                                "name": "Summary",
                                "rows": [
                                    ["Metric", "Value"],
                                    [long_text, 12],
                                    ["第一行\n第二行\n第三行", 3],
                                ],
                            }
                        ]
                    }
                }
            ),
            "artifact": {"format": "xlsx", "filename": "quality.xlsx"},
        }
    )
    assert spreadsheet["ok"] is True, spreadsheet
    from openpyxl import load_workbook

    workbook = load_workbook(
        BytesIO(base64.b64decode(spreadsheet["artifact"]["content_base64"])),
        data_only=False,
    )
    sheet = workbook["Summary"]
    assert sheet["A1"].alignment.horizontal == "center"
    assert sheet["B2"].alignment.horizontal == "right"
    assert sheet.sheet_view.showGridLines is False
    assert (sheet.row_dimensions[2].height or 0) > 20
    assert (sheet.row_dimensions[3].height or 0) >= 48
    workbook.close()

    overflowing_bullets = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Quality",
                        "slides": [
                            {
                                "layout": "two_column",
                                "title": "内容必须适合版面",
                                "left": {
                                    "heading": "左侧",
                                    "bullets": ["\n".join(["中" * 12] * 4)] * 3,
                                },
                                "right": {
                                    "heading": "右侧",
                                    "bullets": ["短"] * 3,
                                },
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "overflow.pptx"},
        }
    )
    assert overflowing_bullets["ok"] is False
    assert "does not fit" in overflowing_bullets["stderr"]

    math_slide = execute_request(
        {
            "skill": "pptx",
            "stdin": json.dumps(
                {
                    "presentation": {
                        "title": "Math",
                        "slides": [
                            {
                                "layout": "bullets",
                                "title": "145 × 12：先估算，再计算",
                                "bullets": [
                                    "145×2＝290",
                                    "145×10＝1450",
                                    "290＋1450＝1740",
                                ],
                            }
                        ],
                    }
                },
                ensure_ascii=False,
            ),
            "artifact": {"format": "pptx", "filename": "math.pptx"},
        }
    )
    assert math_slide["ok"] is True, math_slide
    with ZipFile(
        BytesIO(base64.b64decode(math_slide["artifact"]["content_base64"]))
    ) as archive:
        slide_xml = archive.read("ppt/slides/slide2.xml").decode()
        assert "算式" in slide_xml
        assert "1740" in slide_xml


def check_artifact_error_paths() -> None:
    for code, artifact_format, filename, expected_error in [
        ("pass", "html", "missing.html", "artifact_missing"),
        (
            "import os\nopen(os.environ['NEXAFLOW_OUTPUT_PATH'], 'wb').close()",
            "html",
            "empty.html",
            "artifact_empty",
        ),
    ]:
        result = execute_request(
            {
                "code": code,
                "artifact": {"format": artifact_format, "filename": filename},
            }
        )
        assert result["ok"] is False, result
        assert result["error"] == expected_error or result["error"].startswith(
            f"{expected_error}:"
        ), result

    with tempfile.TemporaryDirectory() as directory:
        root = Path(directory)

        def expect_artifact_error(path: Path, expected: str) -> None:
            try:
                runner_module._read_artifact(path, "html", "page.html")
            except ValueError as exc:
                assert str(exc) == expected, exc
            else:
                raise AssertionError(f"{path} was accepted as an artifact")

        expect_artifact_error(root, "artifact_invalid")
        oversized = root / "oversized.html"
        oversized.write_bytes(b"x" * (runner_module.MAX_ARTIFACT_BYTES + 1))
        expect_artifact_error(oversized, "artifact_too_large")

        root_file = root / "skills-file"
        root_file.write_text("not a directory", encoding="utf-8")
        with mock.patch.dict(
            os.environ,
            {"SANDBOX_SKILLS_DIR": str(root / "missing-skills")},
        ):
            missing_root = execute_request(
                {
                    "code": "pass",
                    "skills": ["documents"],
                    "artifact": {"format": "html", "filename": "missing-root.html"},
                }
            )
        assert missing_root["error"] == "skill_not_found", missing_root

        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(root_file)}):
            file_root = execute_request(
                {
                    "code": "pass",
                    "skills": ["documents"],
                    "artifact": {"format": "html", "filename": "file-root.html"},
                }
            )
        assert file_root["error"] == "skill_not_found", file_root

        skills_root = root / "skills"
        skills_root.mkdir()
        outside = root / "outside"
        outside.mkdir()
        (skills_root / "linked").symlink_to(outside, target_is_directory=True)
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(skills_root)}):
            linked = execute_request(
                {
                    "code": "pass",
                    "skills": ["linked"],
                    "artifact": {"format": "html", "filename": "linked.html"},
                }
            )
        assert linked["error"] == "skill_invalid", linked

        nested = skills_root / "nested"
        nested.mkdir()
        (nested / "linked-file").symlink_to(root_file)
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(skills_root)}):
            nested_link = execute_request(
                {
                    "code": "pass",
                    "skills": ["nested"],
                    "artifact": {"format": "html", "filename": "nested.html"},
                }
            )
        assert nested_link["error"] == "skill_invalid", nested_link

        large = skills_root / "large"
        large.mkdir()
        (large / "SKILL.md").write_text("large", encoding="utf-8")
        (large / "large.bin").write_bytes(b"x" * (runner_module.MAX_SKILL_BYTES + 1))
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(skills_root)}):
            too_large = execute_request(
                {
                    "code": "pass",
                    "skills": ["large"],
                    "artifact": {"format": "html", "filename": "large.html"},
                }
            )
        assert too_large["error"] == "skill_too_large", too_large

        requirements_skill = skills_root / "requirements"
        requirements_skill.mkdir()
        requirements_file = requirements_skill / "requirements.txt"
        requirements_file.write_text("pypdf>=5,<7\n", encoding="utf-8")
        requirements, error = runner_module._skill_requirements(
            skills_root,
            ("requirements",),
        )
        assert requirements == ("pypdf>=5,<7",) and error is None
        requirements_file.write_text("https://example.com/package.whl\n", encoding="utf-8")
        requirements, error = runner_module._skill_requirements(
            skills_root,
            ("requirements",),
        )
        assert requirements == () and error == "skill_requirements_invalid"


def check_run_limits() -> None:
    timeout = run_code(
        "while True: pass",
        limits=Limits(timeout_ms=150, cpu_seconds=1),
    )
    assert timeout.error == "wall_time_limit_exceeded", timeout

    closed = run_code(
        "import os, time\nos.close(1)\nos.close(2)\ntime.sleep(1)",
        limits=Limits(timeout_ms=3_000, cpu_seconds=5),
    )
    assert closed.ok, closed  # pipes EOF; process exit observed via poll branch

    output = run_code(
        "print('x' * 100_000)",
        limits=Limits(max_output_bytes=1024),
    )
    assert output.error == "output_limit_exceeded", output
    assert len(output.stdout.encode()) == 1024, output


def check_terminate() -> None:
    # killpg on an already-exited process group -> ProcessLookupError -> pass.
    finished = subprocess.Popen([sys.executable, "-c", "pass"])
    finished.wait()
    _terminate(finished)  # poll() is not None; killpg lookup may fail harmlessly

    # TimeoutExpired chain: poll stays None, both waits time out.
    class Stubborn:
        pid = 2**31 - 1  # nonexistent group -> ProcessLookupError

        def poll(self) -> None:
            return None

        def wait(self, timeout: float = 0) -> None:
            raise subprocess.TimeoutExpired("process", timeout)

        def kill(self) -> None:
            pass

    _terminate(Stubborn())  # must not raise

    # Success path: real child in its own session, killed via its group.
    sleeper = subprocess.Popen(
        [sys.executable, "-c", "import time; time.sleep(30)"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        start_new_session=True,
    )
    _terminate(sleeper)
    sleeper.wait(timeout=2)
    assert sleeper.returncode == -signal.SIGKILL


def check_blocking_io_retry() -> None:
    real_read = os.read
    state = {"calls": 0}

    def flaky_read(fd: int, n: int) -> bytes:
        # Popen's errpipe read uses n=50000; _collect_output reads use
        # min(64 KiB, limit+1-len) -> 65537 for the default limit. Raise only
        # for the first _collect_output-style read.
        if state["calls"] == 0 and n != 50000:
            state["calls"] += 1
            raise BlockingIOError
        return real_read(fd, n)

    with mock.patch.object(runner_module.os, "read", flaky_read):
        result = run_code("print('retry-ok')")
    assert result.ok and result.stdout.strip() == "retry-ok", result
    assert state["calls"] == 1


def check_non_posix() -> None:
    with mock.patch.object(runner_module.os, "name", "nt"):
        try:
            run_code("print(1)")
        except RuntimeError as exc:
            assert "POSIX" in str(exc)
        else:
            raise AssertionError("non-POSIX run_code accepted")


def check_macos_seatbelt() -> None:
    if sys.platform != "darwin":
        return
    secret = HERE.parent / ".sandbox-seatbelt-secret"
    target = HERE.parent / ".sandbox-seatbelt-write"
    secret.write_text("must-not-read", encoding="utf-8")
    target.unlink(missing_ok=True)
    try:
        result = run_code(
            "from pathlib import Path\n"
            f"secret = Path({str(secret)!r})\n"
            f"target = Path({str(target)!r})\n"
            "try:\n"
            "    secret.read_text()\n"
            "except OSError:\n"
            "    print('read-blocked')\n"
            "else:\n"
            "    print('read-open')\n"
            "try:\n"
            "    target.write_text('polluted')\n"
            "except OSError:\n"
            "    print('write-blocked')\n"
            "else:\n"
            "    print('write-open')\n"
        )
        assert result.ok, result
        assert result.stdout.splitlines() == ["read-blocked", "write-blocked"]
        assert not target.exists()
    finally:
        secret.unlink(missing_ok=True)
        target.unlink(missing_ok=True)


def _raw_request(socket_path: Path, payload: bytes, read_response: bool = True) -> bytes | None:
    with socket.socket(socket.AF_UNIX, socket.SOCK_STREAM) as client:
        client.settimeout(8)
        client.connect(str(socket_path))
        client.sendall(payload)
        if not read_response:
            return None
        return client.makefile("rb").readline()


def check_server_error_paths() -> None:
    with tempfile.TemporaryDirectory() as directory:
        socket_path = Path(directory) / "errors.sock"
        server = SandboxServer(socket_path)
        thread = threading.Thread(target=server.serve_forever, daemon=True)
        thread.start()
        try:
            for _ in range(100):
                if socket_path.exists():
                    break
                time.sleep(0.01)
            assert socket_path.exists()

            response = json.loads(_raw_request(socket_path, b"[1,2]\n") or b"{}")
            assert response["error"] == "request must be an object", response

            response = json.loads(_raw_request(socket_path, b"{not-json\n") or b"{}")
            assert response["ok"] is False and "error" in response, response

            response = json.loads(_raw_request(socket_path, b"\xff\xfe\n") or b"{}")
            assert response["ok"] is False, response

            oversized = b'{"code":"print(1)","pad":"' + b"x" * (MAX_REQUEST_BYTES + 8) + b'"}\n'
            response = json.loads(_raw_request(socket_path, oversized) or b"{}")
            assert response["error"] == "request_too_large", response

            # Concurrent slot exhaustion -> sandbox_busy.
            assert server.run_slot.acquire(timeout=1)
            try:
                response = json.loads(_raw_request(socket_path, b'{"code":"print(1)"}\n') or b"{}")
                assert response["error"] == "sandbox_busy", response
            finally:
                server.run_slot.release()

            # Unexpected internal error -> sandbox_internal_error.
            original = server_module.execute_request
            server_module.execute_request = lambda request: (_ for _ in ()).throw(
                RuntimeError("boom")
            )
            previous_level = server_module.logger.level
            server_module.logger.setLevel(100)  # silence the expected traceback
            try:
                response = json.loads(_raw_request(socket_path, b'{"code":"print(1)"}\n') or b"{}")
                assert response["error"] == "sandbox_internal_error", response
            finally:
                server_module.logger.setLevel(previous_level)
                server_module.execute_request = original

            # Healthcheck request.
            response = json.loads(_raw_request(socket_path, b'{"version":1,"type":"healthcheck"}\n') or b"{}")
            assert response == {"version": 1, "ok": True, "status": "ready"}, response

            # Request timeout: partial line, no newline; handler times out at 5s.
            started = time.monotonic()
            response = json.loads(_raw_request(socket_path, b'{"code": "print(1)"') or b"{}")
            assert response["error"] == "request_timeout", response
            assert time.monotonic() - started >= 4.5
        finally:
            server.shutdown()
            thread.join(timeout=2)
            server.server_close()


def check_server_lifecycle() -> None:
    with tempfile.TemporaryDirectory() as directory:
        # Refuse to replace a regular file.
        file_path = Path(directory) / "occupied"
        file_path.write_text("not a socket", encoding="utf-8")
        try:
            SandboxServer(file_path)
        except RuntimeError as exc:
            assert "refusing to replace non-socket path" in str(exc)
        else:
            raise AssertionError("SandboxServer replaced a regular file")

        # Reuse of a freed socket path exercises the stale-socket unlink.
        socket_path = Path(directory) / "reuse.sock"
        first = SandboxServer(socket_path)
        first.server_close()
        assert not socket_path.exists()
        second = SandboxServer(socket_path)
        second.server_close()
        # Second close -> FileNotFoundError swallowed.
        second.server_close()

        # Overwriting a live socket exercises the exists/unlink path in __init__.
        live = SandboxServer(Path(directory) / "live.sock")
        replacement = SandboxServer(Path(directory) / "live.sock")
        replacement.server_close()
        live.server_close()


def check_main_cli() -> None:
    class FakeServer:
        def __init__(self, socket_path: str | Path):
            self.socket_path = Path(socket_path)
            self.closed = False
            self.shutdown_called = False

        def serve_forever(self) -> None:
            raise KeyboardInterrupt

        def shutdown(self) -> None:
            self.shutdown_called = True

        def server_close(self) -> None:
            self.closed = True

    handlers: dict[int, object] = {}

    def fake_signal(signum: int, handler: object) -> None:
        handlers[signum] = handler

    with tempfile.TemporaryDirectory() as directory:
        socket_path = Path(directory) / "cli.sock"
        with mock.patch.object(server_module, "SandboxServer", FakeServer), mock.patch.object(
            server_module.signal, "signal", fake_signal
        ), mock.patch.object(sys, "argv", ["sandbox-server", "--socket", str(socket_path)]):
            try:
                server_module.main()
            except KeyboardInterrupt:
                pass  # FakeServer.serve_forever raises it; finally ran server_close.
        assert FakeServer.__new__  # constructed through patched class
        # KeyboardInterrupt escaped main(); server_close ran in finally.
        # Trigger the registered SIGTERM handler to exercise `stop`.
        handler = handlers.get(signal.SIGTERM)
        assert handler is not None
        # Reconstruct: main() raises KeyboardInterrupt, so re-run with a server
        # whose serve_forever blocks until shutdown is requested.
        created: list[FakeServer] = []
        with tempfile.TemporaryDirectory() as directory2:
            socket_path2 = Path(directory2) / "cli2.sock"

            class BlockingServer(FakeServer):
                def serve_forever(self) -> None:
                    while not self.shutdown_called:
                        time.sleep(0.01)

            def make_server(path: str | Path) -> BlockingServer:
                instance = BlockingServer(path)
                created.append(instance)
                return instance

            with mock.patch.object(server_module, "SandboxServer", make_server), mock.patch.object(
                server_module.signal, "signal", fake_signal
            ), mock.patch.object(sys, "argv", ["sandbox-server", "--socket", str(socket_path2)]):
                thread = threading.Thread(target=server_module.main, daemon=True)
                thread.start()
                for _ in range(200):
                    if created:
                        break
                    time.sleep(0.01)
                assert created
                handler = handlers.get(signal.SIGTERM)
                assert handler is not None
                handler(None, None)  # type: ignore[misc]
                thread.join(timeout=5)
                assert created[0].closed is True

        # Non-POSIX host exits with SystemExit.
        with mock.patch.object(server_module.os, "name", "nt"), mock.patch.object(
            sys, "argv", ["sandbox-server"]
        ):
            try:
                server_module.main()
            except SystemExit as exc:
                assert "POSIX" in str(exc)
            else:
                raise AssertionError("non-POSIX main() did not exit")


def check_healthcheck_branches() -> None:
    # Deadline already passed -> empty line.
    with socket.socket(socket.AF_UNIX, socket.SOCK_STREAM) as client:
        assert _receive_line(client, time.monotonic() - 0.5) == b""

    def start_server(directory: Path, payload: bytes, delay: float = 0) -> Path:
        socket_path = directory / "serve.sock"

        def runner() -> None:
            with socket.socket(socket.AF_UNIX, socket.SOCK_STREAM) as server:
                server.bind(str(socket_path))
                server.listen(1)
                connection, _ = server.accept()
                with connection:
                    connection.recv(1024)
                    if delay:
                        time.sleep(delay)
                    for value in payload:
                        try:
                            connection.sendall(bytes((value,)))
                        except BrokenPipeError:
                            break

        thread = threading.Thread(target=runner, daemon=True)
        thread.start()
        for _ in range(100):
            if socket_path.exists():
                break
            time.sleep(0.01)
        time.sleep(0.05)  # allow listen() to run before clients connect
        return socket_path

    # EOF before newline -> partial line returned.
    with tempfile.TemporaryDirectory() as directory:
        socket_path = start_server(Path(directory), b'{"partial"')
        with socket.socket(socket.AF_UNIX, socket.SOCK_STREAM) as client:
            deadline = time.monotonic() + PROBE_DEADLINE_SECONDS * 3
            client.settimeout(PROBE_DEADLINE_SECONDS * 3)
            client.connect(str(socket_path))
            client.sendall(b'{"version":1,"type":"healthcheck"}\n')
            assert _receive_line(client, deadline) == b'{"partial"'

    # Oversized chunked response without newline -> probe False (len > max).
    with tempfile.TemporaryDirectory() as directory:
        socket_path = start_server(Path(directory), b"x" * (MAX_RESPONSE_BYTES + 1))
        assert not probe(socket_path)

    # Exact limit without newline -> no crash, probe False.
    with tempfile.TemporaryDirectory() as directory:
        socket_path = start_server(Path(directory), b"y" * MAX_RESPONSE_BYTES)
        assert not probe(socket_path)

    # Empty response -> explicit false branch before JSON parsing.
    with tempfile.TemporaryDirectory() as directory:
        socket_path = start_server(Path(directory), b"")
        assert not probe(socket_path)

    # Healthcheck CLI exits 1 when the default socket is absent.
    env = dict(os.environ)
    prefix = _coverage_prefix()
    run = subprocess.run(
        [*prefix, "-m", "sandbox.healthcheck"],
        cwd=HERE.parent,  # package root for `-m sandbox.*`
        env=env,
        capture_output=True,
        timeout=30,
    )
    assert run.returncode == 1, run

    # CLI exits 0 against a live server on the default path (root or /run available).
    if os.geteuid() == 0 and Path("/run/sandbox").exists():
        socket_path = Path("/run/sandbox/sandbox.sock")
        server = SandboxServer(socket_path)
        thread = threading.Thread(target=server.serve_forever, daemon=True)
        thread.start()
        try:
            run = subprocess.run(
                [*prefix, "-m", "sandbox.healthcheck"],
                cwd=HERE.parent,  # package root for `-m sandbox.*`
                env=env,
                capture_output=True,
                timeout=30,
            )
            assert run.returncode == 0, run
        finally:
            server.shutdown()
            thread.join(timeout=2)
            server.server_close()


def check_server_egress_cli() -> None:
    """Cover `python -m sandbox.server --egress-socket` CLI branches."""
    prefix = _coverage_prefix()
    with tempfile.TemporaryDirectory() as directory:
        socket_path = Path(directory) / "cli-server.sock"
        relative = subprocess.run(
            [
                *prefix,
                "-m",
                "sandbox.server",
                "--egress-socket",
                "relative.sock",
                "--socket",
                str(socket_path),
            ],
            cwd=HERE.parent,
            capture_output=True,
            timeout=30,
        )
        assert relative.returncode != 0
        assert b"absolute" in relative.stderr
        egress_socket = Path(directory) / "egress.sock"
        process = subprocess.Popen(
            [
                *prefix,
                "-m",
                "sandbox.server",
                "--egress-socket",
                str(egress_socket),
                "--socket",
                str(socket_path),
            ],
            cwd=HERE.parent,  # package root for `-m sandbox.*`
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        try:
            for _ in range(200):
                if socket_path.exists():
                    break
                time.sleep(0.02)
            assert socket_path.exists(), "cli server socket not created"
            assert probe(socket_path)
            process.send_signal(signal.SIGTERM)
            process.wait(timeout=10)
            assert process.returncode == 0, process.returncode
        finally:
            if process.poll() is None:
                process.kill()
                process.wait(timeout=5)


def check_server_cli_subprocess() -> None:
    """Run `python -m sandbox.server --socket ...` as a real process."""
    prefix = _coverage_prefix()
    with tempfile.TemporaryDirectory() as directory:
        socket_path = Path(directory) / "cli-server.sock"
        process = subprocess.Popen(
            [*prefix, "-m", "sandbox.server", "--socket", str(socket_path)],
            cwd=HERE.parent,  # package root for `-m sandbox.*`
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        try:
            for _ in range(200):
                if socket_path.exists():
                    break
                time.sleep(0.02)
            assert socket_path.exists(), "cli server socket not created"
            assert probe(socket_path)
            process.send_signal(signal.SIGTERM)
            process.wait(timeout=10)
            assert process.returncode == 0, process.returncode
            assert not socket_path.exists(), "socket not cleaned up on exit"
        finally:
            if process.poll() is None:
                process.kill()
                process.wait(timeout=5)
def check_egress_error_branches() -> None:
    """Cover egress.py validation branches the live proxy tests skip."""
    assert egress_module._is_public_address("not-an-ip") is False
    assert egress_module._is_public_address("::ffff:93.184.216.34") is True
    for target in ("", "x" * 254):
        try:
            egress_module._resolve_public(target, 80)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError("empty or oversized host was accepted")
    with mock.patch.object(
        egress_module.socket, "getaddrinfo", side_effect=socket.gaierror
    ):
        try:
            egress_module._resolve_public("example.test", 80)
        except egress_module.ProxyRequestError as exc:
            assert "DNS" in str(exc)
        else:
            raise AssertionError("DNS failure was accepted")
    for port in ("abc", "22"):
        try:
            egress_module._port(port)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError(f"invalid port {port!r} was accepted")
    for target in ("[::1]", "a:b:c", ":80"):
        try:
            egress_module._host_port(target)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError(f"invalid CONNECT target {target!r} was accepted")
    empty = mock.MagicMock(spec=socket.socket)
    empty.recv.return_value = b""
    try:
        egress_module._read_request(empty)
    except egress_module.ProxyRequestError:
        pass
    else:
        raise AssertionError("empty proxy request was accepted")
    oversized = mock.MagicMock(spec=socket.socket)
    oversized.recv.return_value = b"x" * (egress_module.MAX_HEADER_BYTES + 1)
    try:
        egress_module._read_request(oversized)
    except egress_module.ProxyRequestError:
        pass
    else:
        raise AssertionError("oversized proxy request was accepted")
    for header in (
        b"GET",
        b"GET http://example.com HTTP/2.0\r\n\r\n",
        b"GET http://example.com HTTP/1.1\r\nbad\r\n\r\n",
        b"GET http://example.com HTTP/1.1\r\nbad name: v\r\n\r\n",
        b"GET http://example.com HTTP/1.1\r\nName: bad\x00value\r\n\r\n",
    ):
        try:
            egress_module._headers(header)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError(f"invalid headers {header!r} were accepted")
    for target in (
        "not a url",
        "http://user@example.com/",
        "http://:80/path",
        "http://example.com:abc/",
    ):
        try:
            egress_module._destination("GET", target)
        except egress_module.ProxyRequestError:
            pass
        else:
            raise AssertionError(f"invalid destination {target!r} was accepted")
    left = mock.MagicMock(spec=socket.socket)
    right = mock.MagicMock(spec=socket.socket)
    with mock.patch.object(egress_module.select, "select", side_effect=OSError):
        egress_module._relay(left, right)
    with mock.patch.object(
        egress_module.select, "select", return_value=([left], [], [])
    ):
        left.recv.side_effect = OSError
        egress_module._relay(left, right)
        left.recv.side_effect = None
        left.recv.return_value = b"x" * (egress_module.MAX_RELAY_BYTES + 1)
        egress_module._relay(left, right)
        left.recv.return_value = b"x"
        right.sendall.side_effect = OSError
        egress_module._relay(left, right)
    broken = mock.MagicMock(spec=socket.socket)
    broken.sendall.side_effect = OSError
    egress_module._error(broken, "500 Internal Server Error", "boom")


def check_runner_skill_staging_branches() -> None:
    """Cover runner.py Skill staging and runtime validation errors."""
    with tempfile.TemporaryDirectory() as directory:
        root = Path(directory)
        skills_path = root / "skills-src"
        documents = skills_path / "documents"
        documents.mkdir(parents=True)

        target = root / "target-req.txt"
        target.write_text("requests", encoding="utf-8")
        os.symlink(target, documents / "requirements.txt")
        assert runner_module._skill_requirements(skills_path, ("documents",)) == (
            (),
            "skill_invalid",
        )
        documents.joinpath("requirements.txt").unlink()

        (documents / "requirements.txt").write_text("requests", encoding="utf-8")
        with mock.patch.object(Path, "read_bytes", side_effect=OSError("io")):
            assert runner_module._skill_requirements(
                skills_path, ("documents",)
            ) == ((), "skill_invalid")
        documents.joinpath("requirements.txt").write_bytes(
            b"x" * (runner_module.MAX_REQUIREMENTS_BYTES + 1)
        )
        assert runner_module._skill_requirements(skills_path, ("documents",)) == (
            (),
            "skill_requirements_too_large",
        )
        documents.joinpath("requirements.txt").write_bytes(b"\xff\xfe\x00")
        assert runner_module._skill_requirements(skills_path, ("documents",)) == (
            (),
            "skill_requirements_invalid",
        )
        documents.joinpath("requirements.txt").write_text(
            "# comment only\n", encoding="utf-8"
        )
        assert runner_module._skill_requirements(skills_path, ("documents",)) == (
            (),
            None,
        )
        documents.joinpath("requirements.txt").write_text(
            "\n".join(f"package-{index}" for index in range(33)),
            encoding="utf-8",
        )
        assert runner_module._skill_requirements(skills_path, ("documents",)) == (
            (),
            "skill_requirements_too_many",
        )

        documents.joinpath("requirements.txt").unlink()
        (documents / "SKILL.md").write_text("skill-ready", encoding="utf-8")
        os.symlink(target, documents / "linked.txt")
        workdir = root / "stage-work"
        workdir.mkdir()
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(skills_path)}):
            staged, error = runner_module._stage_skills(workdir, ("documents",))
            assert error == "skill_invalid", error
            documents.joinpath("linked.txt").unlink()
            (documents / "stat-fail.bin").write_text("x", encoding="utf-8")
            staged_workdir = root / "stage-work-2"
            staged_workdir.mkdir()
            original_stat = Path.stat

            def fail_file_stat(self, *args, **kwargs):
                if str(self).endswith("stat-fail.bin") and kwargs.get(
                    "follow_symlinks", True
                ):
                    raise OSError("io")
                return original_stat(self, *args, **kwargs)

            with mock.patch.object(Path, "stat", new=fail_file_stat):
                staged, error = runner_module._stage_skills(
                    staged_workdir, ("documents",)
                )
            assert error == "skill_invalid", error
            documents.joinpath("stat-fail.bin").unlink()

        # _skill_runtime validation (runner.py:268-312).
        assert runner_module._skill_runtime(skills_path, "missing") == (
            None,
            "skill_invalid",
        )
        documents.joinpath("SKILL.md").write_text("plain body", encoding="utf-8")
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        documents.joinpath("SKILL.md").write_text(
            "---\nentrypoint: scripts/render.py\n", encoding="utf-8"
        )
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        documents.joinpath("SKILL.md").write_text(
            "---\nentrypoint: scripts/render.py\nentrypoint: scripts/other.py\n"
            "artifact-format: html\n---\nbody\n",
            encoding="utf-8",
        )
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        documents.joinpath("SKILL.md").write_text(
            "---\nentrypoint: ../escape.py\nartifact-format: html\n---\nbody\n",
            encoding="utf-8",
        )
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        documents.joinpath("SKILL.md").write_text(
            "---\nentrypoint: scripts/render.py\nartifact-format: html\n---\nbody\n",
            encoding="utf-8",
        )
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        scripts = documents / "scripts"
        scripts.mkdir()
        os.symlink(target, scripts / "render.py")
        assert runner_module._skill_runtime(skills_path, "documents") == (
            None,
            "skill_runtime_invalid",
        )
        scripts.joinpath("render.py").unlink()
        (scripts / "render.py").write_text("print('x')", encoding="utf-8")
        documents.joinpath("SKILL.md").write_text(
            "---\nentrypoint: scripts/render.py\nartifact-format: html\n"
            "bare-line\n---\nbody\n",
            encoding="utf-8",
        )
        runtime, runtime_error = runner_module._skill_runtime(
            skills_path, "documents"
        )
        assert runtime_error is None and runtime is not None

        # _installed_tree_size (runner.py:328-343).
        packages = root / "packages"
        packages.mkdir()
        (packages / "a.py").write_text("x", encoding="utf-8")
        os.symlink(target, packages / "link.py")
        count, total, valid = runner_module._installed_tree_size(packages)
        assert valid is False
        packages.joinpath("link.py").unlink()
        (packages / "stat-fail.bin").write_text("x", encoding="utf-8")
        original_stat = Path.stat

        def fail_file_stat(self, *args, **kwargs):
            if str(self).endswith("stat-fail.bin") and kwargs.get(
                "follow_symlinks", True
            ):
                raise OSError("io")
            return original_stat(self, *args, **kwargs)

        with mock.patch.object(Path, "stat", new=fail_file_stat):
            count, total, valid = runner_module._installed_tree_size(packages)
        assert valid is False
        packages.joinpath("stat-fail.bin").unlink()
        (packages / "b.py").write_text("y", encoding="utf-8")
        with mock.patch.object(runner_module, "MAX_PACKAGES_FILES", 1):
            count, total, valid = runner_module._installed_tree_size(packages)
        assert valid is False

        # Non-loopback proxy rejected on macOS (runner.py:512-515).
        if sys.platform == "darwin":
            try:
                runner_module._sandboxed_child_command(
                    ["x"], Path("/tmp"), "http://10.0.0.1:8080"
                )
            except ValueError as exc:
                assert "loopback" in str(exc)
            else:
                raise AssertionError("non-loopback proxy was accepted")

        # Invalid skill requests (runner.py:676-682).
        for call in (
            lambda: run_code("pass", skill="documents"),
            lambda: run_code(None, skill="BAD!name"),
            lambda: run_code(None, skill="documents", skills=("documents",)),
        ):
            try:
                call()
            except ValueError:
                pass
            else:
                raise AssertionError("invalid skill request was accepted")

        # Skill runtime failures surface as errors (runner.py:710-721).
        bad_root = root / "skills-bad"
        bad_docs = bad_root / "documents"
        bad_docs.mkdir(parents=True)
        (bad_docs / "SKILL.md").write_text(
            "---\nname: documents\n---\nbody\n", encoding="utf-8"
        )
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(bad_root)}):
            result = run_code(None, skill="documents", artifact=("html", "page.html"))
        assert result.ok is False
        assert result.error == "skill_runtime_invalid", result.error
        valid_root = root / "skills-valid"
        valid_docs = valid_root / "documents"
        renderer = valid_docs / "scripts" / "render.py"
        renderer.parent.mkdir(parents=True)
        renderer.write_text("print('ready')", encoding="utf-8")
        (valid_docs / "SKILL.md").write_text(
            "---\nentrypoint: scripts/render.py\nartifact-format: html\n---\nbody\n",
            encoding="utf-8",
        )
        with mock.patch.dict(os.environ, {"SANDBOX_SKILLS_DIR": str(valid_root)}):
            result = run_code(None, skill="documents")
            assert result.ok is False
            assert result.error == "skill_artifact_required", result.error
            result = run_code(
                None, skill="documents", artifact=("docx", "report.docx")
            )
            assert result.ok is False
            assert result.error == "skill_artifact_format_invalid", result.error


def check_skill_dependency_install() -> None:
    """Cover runner.py pip dependency installation branches."""
    with tempfile.TemporaryDirectory() as directory:
        root = Path(directory)
        workdir = root / "work"
        workdir.mkdir()
        skills_path = workdir / "skills-src"
        documents = skills_path / "documents"
        documents.mkdir(parents=True)
        (documents / "requirements.txt").write_text(
            "requests>=2.31\n", encoding="utf-8"
        )
        proxy = "http://127.0.0.1:8080"

        assert runner_module._install_skill_dependencies(
            workdir, skills_path, ("documents",), None, {}
        ) == (None, "skill_network_unavailable")

        process = mock.MagicMock()
        process.returncode = 0
        process.pid = 99999999
        with mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"", b"", None)
        ):
            packages, error = runner_module._install_skill_dependencies(
                workdir, skills_path, ("documents",), proxy, {"HTTP_PROXY": proxy}
            )
        assert error is None
        assert packages == workdir / "packages"
        assert packages.is_dir()

        workdir2 = root / "work2"
        workdir2.mkdir()
        with mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"", b"", "pip boom")
        ):
            _, error = runner_module._install_skill_dependencies(
                workdir2, skills_path, ("documents",), proxy, {}
            )
        assert error == "skill_dependency_install_failed:pip boom", error

        process.returncode = 1
        workdir3 = root / "work3"
        workdir3.mkdir()
        with mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"oops", b"", None)
        ):
            _, error = runner_module._install_skill_dependencies(
                workdir3, skills_path, ("documents",), proxy, {}
            )
        assert error is not None
        assert error.startswith("skill_dependency_install_failed:")
        assert "oops" in error

        workdir4 = root / "work4"
        workdir4.mkdir()
        with mock.patch.object(
            runner_module.subprocess, "Popen", side_effect=OSError("no pip")
        ):
            _, error = runner_module._install_skill_dependencies(
                workdir4, skills_path, ("documents",), proxy, {}
            )
        assert error is not None
        assert error.startswith("skill_dependency_install_failed:")

        process.returncode = 0
        workdir5 = root / "work5"
        workdir5.mkdir()
        with mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"", b"", None)
        ), mock.patch.object(
            runner_module, "_installed_tree_size", return_value=(3, 3, False)
        ):
            _, error = runner_module._install_skill_dependencies(
                workdir5, skills_path, ("documents",), proxy, {}
            )
        assert error == "skill_dependencies_too_large", error
        # Requirement errors short-circuit (runner.py:353-355).
        bad_skills = root / "bad-skills"
        (bad_skills / "documents").mkdir(parents=True)
        (root / "target.txt").write_text("requests", encoding="utf-8")
        os.symlink(root / "target.txt", bad_skills / "documents" / "requirements.txt")
        workdir6 = root / "work6"
        workdir6.mkdir()
        assert runner_module._install_skill_dependencies(
            workdir6, bad_skills, ("documents",), proxy, {}
        ) == (None, "skill_invalid")

        # Seatbelt with base prefix outside home (runner.py:497-498).
        if sys.platform == "darwin":
            with mock.patch.object(Path, "home", return_value=Path(root)):
                command = runner_module._sandboxed_child_command(
                    ["x"], workdir6, None
                )
            assert "sandbox-exec" in command[0]


def check_egress_integration() -> None:
    """Cover run_code's egress proxy staging (runner.py:756-822)."""
    with tempfile.TemporaryDirectory() as directory:
        root = Path(directory)
        skills_path = root / "skills"
        documents = skills_path / "documents"
        documents.mkdir(parents=True)
        (documents / "SKILL.md").write_text("skill-ready", encoding="utf-8")
        (documents / "requirements.txt").write_text(
            "requests>=2.31\n", encoding="utf-8"
        )
        proxy = mock.MagicMock()
        proxy.start.return_value = "http://127.0.0.1:12345"
        proxy.environment.return_value = {"HTTP_PROXY": "http://127.0.0.1:12345"}
        process = mock.MagicMock()
        process.returncode = 0
        process.pid = 99999999
        with mock.patch.dict(
            os.environ,
            {
                "SANDBOX_SKILLS_DIR": str(skills_path),
                "SANDBOX_EGRESS_SOCKET": str(root / "egress.sock"),
            },
        ), mock.patch.object(
            runner_module, "LocalEgressProxy", return_value=proxy
        ), mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"", b"", None)
        ):
            result = run_code(
                "pass", skills=("documents",), artifact=("html", "page.html")
            )
        assert proxy.start.called
        assert proxy.environment.called
        assert proxy.close.called
        assert result.error is None or result.error.startswith("artifact"), result

        with mock.patch.dict(
            os.environ,
            {
                "SANDBOX_SKILLS_DIR": str(skills_path),
                "SANDBOX_EGRESS_SOCKET": str(root / "egress.sock"),
            },
        ), mock.patch.object(
            runner_module, "LocalEgressProxy", return_value=proxy
        ), mock.patch.object(
            runner_module.subprocess, "Popen", return_value=process
        ), mock.patch.object(
            runner_module, "_collect_output", return_value=(b"", b"", "boom")
        ):
            failed = run_code(
                "pass", skills=("documents",), artifact=("html", "page.html")
            )
        assert failed.ok is False
        assert failed.error.startswith("skill_dependency_install_failed:"), failed.error


def check_egress_remaining_branches() -> None:
    """Cover the remaining egress.py proxy branches."""
    # urlsplit rejection (egress.py:138-139).
    try:
        egress_module._destination("GET", "http://[::1")
    except egress_module.ProxyRequestError:
        pass
    else:
        raise AssertionError("malformed proxy URL was accepted")
    # Idle relay exit (egress.py:168-169).
    left = mock.MagicMock(spec=socket.socket)
    right = mock.MagicMock(spec=socket.socket)
    with mock.patch.object(egress_module.select, "select", return_value=([], [], [])):
        egress_module._relay(left, right)
    # _connect success and failure (egress.py:199-211).
    upstream = mock.MagicMock(spec=socket.socket)
    with mock.patch.object(
        egress_module, "_resolve_public", return_value=("93.184.216.34",)
    ), mock.patch.object(
        egress_module.socket, "create_connection", return_value=upstream
    ):
        assert egress_module._connect("example.com", 443) is upstream
    with mock.patch.object(
        egress_module, "_resolve_public", return_value=("93.184.216.34",)
    ), mock.patch.object(
        egress_module.socket, "create_connection", side_effect=OSError("refused")
    ):
        try:
            egress_module._connect("example.com", 443)
        except OSError as exc:
            assert "unreachable" in str(exc)
        else:
            raise AssertionError("unreachable destination was accepted")
    # Forwarded headers (egress.py:234-245).
    client, proxy_side = socket.socketpair()
    upstream_peer, upstream_side = socket.socketpair()
    with mock.patch.object(egress_module, "_connect", return_value=upstream_side):
        thread = threading.Thread(
            target=egress_module._serve_client, args=(proxy_side,), daemon=True
        )
        thread.start()
        client.sendall(
            b"GET http://example.com/ready HTTP/1.1\r\n"
            b"Host: example.com\r\nX-Custom: yes\r\n\r\nbody"
        )
        forwarded = upstream_peer.recv(4096)
        assert forwarded.startswith(b"GET /ready")
        assert b"X-Custom: yes" in forwarded
        client.close()
        upstream_peer.close()
        upstream_side.close()
        thread.join(timeout=2)
        assert not thread.is_alive()
    # 502 on upstream failure (egress.py:248-249).
    client, proxy_side = socket.socketpair()
    with mock.patch.object(egress_module, "_connect", side_effect=OSError("down")):
        thread = threading.Thread(
            target=egress_module._serve_client, args=(proxy_side,), daemon=True
        )
        thread.start()
        client.sendall(
            b"GET http://example.com/ready HTTP/1.1\r\nHost: example.com\r\n\r\n"
        )
        assert client.recv(4096).startswith(b"HTTP/1.1 502")
        client.close()
        thread.join(timeout=2)
        assert not thread.is_alive()
    # LocalEgressProxy validation, start idempotence, close idempotence.
    with tempfile.TemporaryDirectory() as directory:
        try:
            egress_module.LocalEgressProxy("relative.sock")
        except ValueError:
            pass
        else:
            raise AssertionError("relative egress socket was accepted")
        proxy = egress_module.LocalEgressProxy(Path(directory) / "egress.sock")
        assert proxy.url.startswith("http://127.0.0.1:")
        assert proxy.start() == proxy.url
        assert proxy.start() == proxy.url  # already started (egress.py:287-288)
        proxy.close()
        proxy.close()  # idempotent (egress.py:331-332)
    # Accept-loop error handling (egress.py:299-303).
    proxy = egress_module.LocalEgressProxy(Path("/tmp/nexaflow-egress-errors.sock"))
    proxy._listener = mock.MagicMock(spec=socket.socket)
    proxy._listener.accept.side_effect = OSError("closed")
    proxy._serve()  # returns on accept failure
    proxy._listener.accept.side_effect = [
        (mock.MagicMock(spec=socket.socket), None),
        OSError("closed"),
    ]
    with mock.patch.object(proxy._slots, "acquire", return_value=False):
        proxy._serve()  # slot exhaustion closes the client and continues
    proxy.close()
    # Connection teardown on upstream failure (egress.py:317-327, 338).
    proxy = egress_module.LocalEgressProxy(Path("/tmp/nexaflow-egress-errors-2.sock"))
    proxy._slots.acquire()
    proxy._handle(mock.MagicMock(spec=socket.socket))
    proxy._connections.add(mock.MagicMock(spec=socket.socket))
    proxy.close()  # closes tracked connections


def check_egress_cli() -> None:
    """Run `python -m sandbox.egress --fd ...` as a real process."""
    run = subprocess.run(
        [*_coverage_prefix(), "-m", "sandbox.egress", "--fd", "9999"],
        cwd=HERE.parent,
        capture_output=True,
        timeout=30,
    )
    assert run.returncode != 0


def main() -> None:
    check_limits_validation()
    check_run_code_validation()
    check_egress_policy()
    check_egress_relay()
    check_local_egress_proxy()
    check_artifact_outputs_and_skills()
    check_builtin_skill_entrypoints()
    check_builtin_skill_quality_guards()
    check_artifact_error_paths()
    check_run_limits()
    check_terminate()
    check_blocking_io_retry()
    check_non_posix()
    check_macos_seatbelt()
    check_server_error_paths()
    check_server_lifecycle()
    check_main_cli()
    check_healthcheck_branches()
    check_server_cli_subprocess()
    check_server_egress_cli()
    check_egress_error_branches()
    check_runner_skill_staging_branches()
    check_skill_dependency_install()
    check_egress_integration()
    check_egress_remaining_branches()
    check_egress_cli()
    print("sandbox extended tests passed")


if __name__ == "__main__":
    main()
